import os
from pathlib import Path
from typing import List, Tuple, Dict
import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from models import ImageData, TextData, Config
from image_filter import ImageFilter
from composite_detector import CompositeChartDetector
import json


class DocumentExtractor:
    def __init__(self, config_path: str = "config.json"):
        with open(config_path, 'r') as f:
            config_data = json.load(f)
        self.config = Config(**config_data)
        self.config_path = config_path
        self._setup_directories()
        
        # Inicializar filtro de imágenes si está habilitado
        filter_config = self.config.extraction.get('image_filter', {})
        self.filter_enabled = filter_config.get('enabled', True)
        if self.filter_enabled:
            self.image_filter = ImageFilter(config_path)
        else:
            self.image_filter = None
            print("⚠️  Filtro de imágenes deshabilitado")
        
        # Inicializar detector de gráficos compuestos
        composite_config = self.config.extraction.get('composite_detection', {})
        self.composite_enabled = composite_config.get('enabled', True)
        if self.composite_enabled:
            self.composite_detector = CompositeChartDetector(config_path)
        else:
            self.composite_detector = None
            print("⚠️  Detección de gráficos compuestos deshabilitada")
    
    def _setup_directories(self):
        """Crea los directorios necesarios"""
        base_dir = Path(self.config.extraction['output_dir'])
        for subdir in ['images_dir', 'text_dir', 'data_dir']:
            path = base_dir / self.config.extraction[subdir]
            path.mkdir(parents=True, exist_ok=True)
    
    def extract_pdf(self, pdf_path: str) -> Tuple[List[TextData], List[ImageData]]:
        """Extrae texto e imágenes de un PDF"""
        doc = fitz.open(pdf_path)
        text_data = []
        image_data = []
        
        for page_num, page in enumerate(doc, start=1):
            # Extraer texto
            text = page.get_text()
            if text.strip():
                text_data.append(TextData(
                    page_number=page_num,
                    content=text
                ))
            
            # Extraer imágenes
            images = page.get_images()
            for img_idx, img in enumerate(images):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                
                # Guardar imagen
                filename = f"page{page_num}_img{img_idx}.{image_ext}"
                img_path = Path(self.config.extraction['output_dir']) / \
                          self.config.extraction['images_dir'] / filename
                
                with open(img_path, "wb") as f:
                    f.write(image_bytes)
                
                # Obtener dimensiones
                img_obj = Image.open(img_path)
                w, h = img_obj.size
                
                image_data.append(ImageData(
                    filename=filename,
                    page_number=page_num,
                    path=str(img_path),
                    width=w,
                    height=h
                ))
        
        doc.close()
        return text_data, image_data
    
    def extract_pptx(self, pptx_path: str) -> Tuple[List[TextData], List[ImageData]]:
        """Extrae texto e imágenes de un PPTX"""
        prs = Presentation(pptx_path)
        text_data = []
        image_data = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            # Extraer texto
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
            
            if text_content:
                text_data.append(TextData(
                    page_number=slide_num,
                    content="\n".join(text_content)
                ))
            
            # Extraer imágenes
            for img_idx, shape in enumerate(slide.shapes):
                if shape.shape_type == 13:  # Picture
                    image = shape.image
                    image_bytes = image.blob
                    
                    filename = f"slide{slide_num}_img{img_idx}.{image.ext}"
                    img_path = Path(self.config.extraction['output_dir']) / \
                              self.config.extraction['images_dir'] / filename
                    
                    with open(img_path, "wb") as f:
                        f.write(image_bytes)
                    
                    img_obj = Image.open(img_path)
                    w, h = img_obj.size
                    
                    image_data.append(ImageData(
                        filename=filename,
                        page_number=slide_num,
                        path=str(img_path),
                        width=w,
                        height=h
                    ))
        
        return text_data, image_data
    
    def extract(self, file_path: str) -> Tuple[List[TextData], List[ImageData]]:
        """Extrae contenido según el tipo de archivo"""
        ext = Path(file_path).suffix.lower()
        
        if ext == '.pdf':
            text_data, image_data = self.extract_pdf(file_path)
        elif ext in ['.ppt', '.pptx']:
            text_data, image_data = self.extract_pptx(file_path)
        else:
            raise ValueError(f"Formato no soportado: {ext}")
        
        # Diccionario para almacenar resultados de OCR (útil para detección de composites)
        ocr_results: Dict[str, int] = {}
        
        # Aplicar filtro de imágenes si está habilitado
        if self.filter_enabled and self.image_filter and image_data:
            valuable_images, discarded_images = self.image_filter.filter_images(image_data)
            
            # Guardar resultados de OCR para las imágenes valiosas
            # (el filtro ya hizo OCR, aprovechamos esa información)
            ocr_results = self.image_filter.get_ocr_results()
            
            # Opcionalmente eliminar archivos de imágenes descartadas
            for img in discarded_images:
                try:
                    Path(img.path).unlink()
                except Exception as e:
                    print(f"  ⚠️  No se pudo eliminar {img.filename}: {e}")
            
            image_data = valuable_images
        
        # Detectar y enriquecer gráficos compuestos (solo para PDFs)
        if self.composite_enabled and self.composite_detector and image_data and ext == '.pdf':
            image_data = self.composite_detector.enrich_images_with_context(
                image_data, 
                file_path,
                ocr_results
            )
        
        return text_data, image_data


if __name__ == "__main__":
    # Ejemplo de uso
    extractor = DocumentExtractor()
    text, images = extractor.extract("document.pdf")
    print(f"Extraídas {len(text)} páginas de texto y {len(images)} imágenes")
